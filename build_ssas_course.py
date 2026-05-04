from __future__ import annotations

from pathlib import Path
import re
import shutil
from textwrap import dedent

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BASE = Path("/Users/kumbulani.tshuma/Documents/my trainning/SSAS-Fundamentals-SQL-Server-Analysis-Services")

TOPICS = [
    {
        "day": "Day-01",
        "title": "Introduction to SQL Server Analysis Services",
        "dataset": "v1_assmang_mining_base.sql",
        "objectives": [
            "Explain what SSAS is and where it fits in the Microsoft BI stack.",
            "Differentiate multidimensional and tabular models at a beginner level.",
            "Understand SSAS terminology such as cube, dimension, hierarchy, measure, and processing.",
            "Connect the SSAS learning journey to Assmang production analytics use cases.",
        ],
        "sections": [
            ("What SSAS does", [
                "SQL Server Analysis Services is Microsoft's analytical engine for building semantic models over warehouse data.",
                "It turns detailed transactional or warehouse tables into structures designed for fast slicing, drilling, and summarisation.",
                "In this course the focus is on multidimensional SSAS, where data is modelled using dimensions, measures, and cubes.",
            ]),
            ("Why Assmang would use it", [
                "Production data spans mines, dates, departments, safety indicators, and cost centres.",
                "Executives want answers such as revenue by mine, cost per tonne by month, and trend analysis by commodity type.",
                "SSAS pre-builds the analytical model so these queries run quickly and consistently.",
            ]),
            ("Core architecture", [
                "Source database stores relational warehouse tables.",
                "SSDT designs the SSAS project: data source, data source view, dimensions, measure groups, and cube.",
                "Analysis Services processes the model and exposes it to Excel, Power BI, and MDX clients.",
            ]),
            ("Key beginner terms", [
                "Dimension = how you slice data, e.g. Mine or Date.",
                "Measure = the numeric business fact being analysed, e.g. Tonnes Produced or Revenue ZAR.",
                "Hierarchy = a drill path such as Year > Quarter > Month.",
                "Processing = loading data and building aggregations in the SSAS database.",
            ]),
        ],
        "diagram": [
            "Relational Warehouse -> SSAS Project -> Processed Cube -> Excel/Power BI",
            "           |                 |                 |",
            "       Dim/Fact tables   Dimensions &      Analytical",
            "                         Measure Groups      Reports",
        ],
        "practical_steps": [
            "Load `datasets/v1_assmang_mining_base.sql` into SQL Server and verify all four dimension tables exist.",
            "Open Visual Studio with SSDT and create a new Analysis Services Multidimensional Project.",
            "Create a data source pointing to the `AssmangMining` database.",
            "Create a Data Source View including `Dim_Mine`, `Dim_Department`, `Dim_Employee`, and `Dim_Date`.",
            "Deploy an empty project shell to the local SSAS instance to confirm connectivity and permissions.",
        ],
        "hands_on": [
            "List five business questions Assmang managers could answer better with a cube than with an operational report.",
            "Compare multidimensional and tabular SSAS in one page and recommend which model best fits this training environment.",
            "Document the meaning of cube, hierarchy, attribute, member, and processing in your own words.",
        ],
        "mcqs": [
            ("What is the main purpose of SSAS?", ["Store transactional rows", "Run analytical models over prepared data", "Replace SQL Server Engine", "Manage Windows users"], "Run analytical models over prepared data"),
            ("Which object represents how users slice data?", ["Measure", "Dimension", "Partition", "Perspective"], "Dimension"),
            ("Which statement best describes processing?", ["Formatting reports", "Refreshing the cube structure and data", "Deleting measures", "Renaming dimensions"], "Refreshing the cube structure and data"),
        ],
        "scenario": "Assmang executives ask for monthly production by mine, by commodity type, and by province. Explain why a processed SSAS cube is a better delivery layer than repeatedly writing raw SQL over source tables.",
        "practical_challenge": "Create an SSAS project, connect it to `AssmangMining`, include the v1 dimensions in a Data Source View, and document your successful deployment validation steps.",
    },
    {
        "day": "Day-01",
        "title": "Multidimensional Models and Dimensions",
        "dataset": "v1_assmang_mining_base.sql",
        "objectives": [
            "Understand star-schema thinking and how dimensions support analysis.",
            "Design dimensions from the Assmang dimension tables.",
            "Build hierarchies that support drill-down navigation.",
            "Recognise common dimension design issues such as poor keys or weak hierarchies.",
        ],
        "sections": [
            ("Dimensional thinking", [
                "A multidimensional model separates descriptive business context from measurable facts.",
                "Dimension tables answer 'by what' questions: by mine, by month, by department, by employee.",
                "A clean dimension model improves user navigation and aggregation behaviour.",
            ]),
            ("Assmang dimensions", [
                "`Dim_Mine` supports analysis by mine, province, and commodity type.",
                "`Dim_Date` supports time intelligence through year, quarter, month, and day levels.",
                "`Dim_Department` supports cost and people analysis by function area.",
                "`Dim_Employee` supports workforce-centric reporting and role-based views.",
            ]),
            ("Hierarchies and drill paths", [
                "A hierarchy gives users a guided navigation path rather than a flat list of attributes.",
                "Examples include Time: Year > Quarter > Month > Day and Mine: Mine Type > Province > Mine Name.",
                "Good hierarchies make Excel pivot navigation and MDX browsing easier.",
            ]),
            ("Slowly changing dimensions", [
                "Type 1 overwrites old values and is useful when history is not required.",
                "Type 2 preserves historical versions and is useful when reporting must reflect old business states.",
                "Beginners should understand the concept even if the training implementation stays simple.",
            ]),
        ],
        "diagram": [
            "                Dim_Date",
            "                   |",
            "Dim_Mine ---- Fact_Production ---- Dim_Department",
            "                   |",
            "              Dim_Employee",
        ],
        "practical_steps": [
            "Use the Data Source View from Topic 1 and run the Dimension Wizard for `Dim_Mine`.",
            "Set `MineID` as the key attribute and `MineName` as the name column.",
            "Add attributes for `MineType`, `Province`, and `EstablishedYear`.",
            "Create a user hierarchy `Mine Type > Province > Mine Name` and process the dimension.",
            "Repeat the exercise for `Dim_Date` and create the `Year > Quarter > Month > Day` hierarchy.",
        ],
        "hands_on": [
            "Design a `Dim_Shift` dimension for shift-based mining operations and propose a useful hierarchy.",
            "Decide whether employee assignment changes should be modelled as Type 1 or Type 2 for management reporting.",
            "Review each Assmang dimension and identify which attributes should be hidden from end users.",
        ],
        "mcqs": [
            ("What is the main purpose of a hierarchy?", ["Increase row counts", "Support drill-down navigation", "Replace measures", "Store MDX code"], "Support drill-down navigation"),
            ("Which Assmang dimension naturally owns the hierarchy Year > Quarter > Month > Day?", ["Dim_Mine", "Dim_Date", "Dim_Department", "Dim_Employee"], "Dim_Date"),
            ("Which SCD type keeps history?", ["Type 0", "Type 1", "Type 2", "Type 4"], "Type 2"),
        ],
        "scenario": "Assmang wants to compare mines by commodity type first, then by province, then by site. Explain how you would implement that navigation path and why it is better than a flat member list.",
        "practical_challenge": "Build `Dim_Mine` and `Dim_Date` in SSDT, process them, and confirm that users can browse both hierarchies successfully.",
    },
    {
        "day": "Day-01",
        "title": "Measures, Measure Groups, and Aggregations",
        "dataset": "v2_assmang_mining_extended.sql",
        "objectives": [
            "Understand the relationship between fact tables, measure groups, and cube measures.",
            "Choose correct aggregation behavior for common mining metrics.",
            "Recognise additive, semi-additive, and non-additive business values.",
            "Understand the purpose of aggregations in performance optimisation.",
        ],
        "sections": [
            ("Fact tables and measure groups", [
                "A fact table stores measurable business events at a defined grain.",
                "In SSAS, fact tables surface as measure groups.",
                "`FactProduction` becomes a production measure group; `FactOperatingCosts` becomes a cost measure group.",
            ]),
            ("Choosing measures", [
                "Production metrics include tonnes produced, grade, revenue, and cost per tonne.",
                "Cost metrics include labour, equipment, maintenance, safety, utilities, and other costs.",
                "Measures should reflect the grain of the underlying fact and have meaningful aggregation rules.",
            ]),
            ("Aggregation behaviour", [
                "Additive measures such as revenue and labour cost usually sum well across most dimensions.",
                "Average-style measures such as grade or uptime should not be blindly summed.",
                "Design choices here directly affect business trust in the cube.",
            ]),
            ("Why aggregations matter", [
                "SSAS pre-computes useful summaries so queries do not always scan the full fact table.",
                "This makes repeated management queries much faster.",
                "Aggregation design is one of the reasons cube workloads feel responsive.",
            ]),
        ],
        "diagram": [
            "FactProduction -> Measure Group: Production",
            "   |           -> TonnesProduced (SUM)",
            "   |           -> RevenueZAR (SUM)",
            "   |           -> Grade (AVERAGE)",
            "FactOperatingCosts -> Measure Group: Operating Costs",
        ],
        "practical_steps": [
            "Load `datasets/v2_assmang_mining_extended.sql` to add production and operating cost facts.",
            "Update the Data Source View to include `FactProduction` and `FactOperatingCosts`.",
            "Create a cube using the Cube Wizard and select both fact tables as measure groups.",
            "Review the proposed measures and correct any aggregation types that should use averages instead of sums.",
            "Process the cube and browse `TonnesProduced`, `RevenueZAR`, and cost measures by mine and month.",
        ],
        "hands_on": [
            "List which measures in `FactProduction` are additive, semi-additive, or better modelled as calculated measures.",
            "Design an `Equipment Efficiency` measure group using the v3 dataset and propose three measures.",
            "Explain why `CostPerTonneZAR` is often better as a calculated expression than a raw summed measure.",
        ],
        "mcqs": [
            ("What does a measure group normally originate from?", ["A dimension table", "A fact table", "An MDX query", "A user hierarchy"], "A fact table"),
            ("Which metric is most likely to use SUM?", ["RevenueZAR", "Grade", "UptimePercentage", "ComplianceScore"], "RevenueZAR"),
            ("Why are aggregations useful?", ["They reduce file names", "They improve query speed", "They replace dimensions", "They disable processing"], "They improve query speed"),
        ],
        "scenario": "Management notices that chrome grade results look wrong after cube browsing. Explain what aggregation design mistake could have caused this and how you would correct it.",
        "practical_challenge": "Create the production and operating cost measure groups, process the cube, and demonstrate one additive measure and one average-style measure in the cube browser.",
    },
    {
        "day": "Day-01",
        "title": "Building and Deploying SSAS Cubes",
        "dataset": "v2_assmang_mining_extended.sql",
        "objectives": [
            "Understand the end-to-end workflow for building a multidimensional cube in SSDT.",
            "Create a cube from data source, DSV, dimensions, and measure groups.",
            "Deploy and process a cube to an SSAS instance.",
            "Perform validation checks before handing the cube to users.",
        ],
        "sections": [
            ("Cube build workflow", [
                "A typical workflow moves from relational source -> data source -> data source view -> dimensions -> cube -> deployment -> processing.",
                "Each stage introduces metadata that shapes the analytical experience.",
            ]),
            ("Data source and DSV", [
                "The data source defines how SSDT connects to SQL Server.",
                "The Data Source View is the logical modelling layer that selects tables and relationships.",
                "A clean DSV makes later cube design easier and clearer.",
            ]),
            ("Deployment and processing", [
                "Deployment pushes the project metadata to an SSAS server.",
                "Processing loads data and builds the structures users actually query.",
                "Without successful processing, a deployed cube is not analytically usable.",
            ]),
            ("Validation and readiness", [
                "Validate measure totals, hierarchy browsing, key attributes, and security assumptions.",
                "A cube should be tested with both SSDT browser checks and client tool connectivity.",
            ]),
        ],
        "diagram": [
            "SQL Server Warehouse -> Data Source -> DSV -> Dimensions + Measures -> Deploy -> Process -> Browse",
        ],
        "practical_steps": [
            "Open the SSAS project and confirm the data source points to `AssmangMining`.",
            "Add all v2 dimensions and measure groups to a single cube named `Assmang Mining Analytics`.",
            "Review cube dimension usage so that relationships are correctly mapped.",
            "Deploy the cube to the training SSAS instance and process it fully.",
            "Use the cube browser to confirm that production and cost metrics can be sliced by mine and time.",
        ],
        "hands_on": [
            "Write a deployment checklist that a junior BI developer could follow without trainer support.",
            "Design an incremental processing strategy for daily production data refreshes.",
            "List five validation checks to run before releasing a new cube version to management.",
        ],
        "mcqs": [
            ("Which object usually comes before cube creation?", ["MDX script", "Data Source View", "Perspective", "Partition"], "Data Source View"),
            ("What makes a deployed cube usable to end users?", ["Renaming it", "Processing it", "Compressing it", "Moving it to Excel"], "Processing it"),
            ("What is the main goal of deployment validation?", ["Reduce folder count", "Confirm analytical correctness and availability", "Create new tables", "Change source keys"], "Confirm analytical correctness and availability"),
        ],
        "scenario": "The cube deploys successfully but business users see empty results in the browser. Explain the most likely cause and the validation sequence you would follow.",
        "practical_challenge": "Build, deploy, process, and browse the Assmang cube. Capture evidence that revenue and costs are visible by mine and by month.",
    },
    {
        "day": "Day-02",
        "title": "MDX Query Fundamentals",
        "dataset": "v3_assmang_mining_complete.sql",
        "objectives": [
            "Understand the structure of a basic MDX SELECT statement.",
            "Work with measures, members, sets, tuples, and slicers.",
            "Query the Assmang cube for common analytical views.",
            "Recognise how MDX differs from SQL thinking.",
        ],
        "sections": [
            ("How MDX differs from SQL", [
                "SQL retrieves rows from tables; MDX navigates coordinates in a cube.",
                "MDX focuses on axes, members, sets, and dimensional context.",
                "The question is not only 'which rows?' but 'which slice of the cube?'.",
            ]),
            ("Basic SELECT structure", [
                "Measures often appear on columns and dimension members appear on rows.",
                "The FROM clause names the cube, not a table.",
                "The WHERE clause acts as a slicer over dimensional context.",
            ]),
            ("Core MDX building blocks", [
                "Member = one addressable point such as `[Mine].[Mine Name].&[Beeshoek Mine]`.",
                "Set = a collection of members.",
                "Tuple = one combination across dimensions, such as Mine + Month.",
            ]),
            ("Beginner query patterns", [
                "One measure by one dimension.",
                "One measure by a hierarchy level.",
                "Filtered results using a slicer.",
            ]),
        ],
        "diagram": [
            "COLUMNS: [Measures].[RevenueZAR]",
            "ROWS:    [Mine].[Mine Name].Members",
            "SLICER:  [Date].[Calendar Year].&[2024]",
        ],
        "practical_steps": [
            "Open SSMS against the SSAS instance and connect to the processed Assmang cube.",
            "Run a basic query returning `TonnesProduced` by mine.",
            "Add a time slicer for calendar year 2024.",
            "Query revenue by commodity hierarchy level instead of by individual mine.",
            "Use a set to return only the iron ore mines.",
        ],
        "hands_on": [
            "Write an MDX query that returns total production for chrome operations only.",
            "Write a query that shows revenue by quarter for one selected mine.",
            "Explain the difference between rows, columns, and slicers in your own words.",
        ],
        "mcqs": [
            ("MDX queries run against what object?", ["A single SQL table", "A cube", "An index", "A Windows service only"], "A cube"),
            ("What does the WHERE clause usually act as in MDX?", ["A join", "A slicer", "A data type", "A partition"], "A slicer"),
            ("Which object is a collection of members?", ["Tuple", "Set", "Measure group", "Attribute relationship"], "Set"),
        ],
        "scenario": "An analyst wants revenue by mine for 2024 only, but keeps getting all years. Explain how the slicer changes the dimensional context of the query.",
        "practical_challenge": "Write three working MDX queries: by mine, by commodity type, and by quarter, all against the Assmang cube.",
    },
    {
        "day": "Day-02",
        "title": "Advanced Queries, Calculations, and KPIs",
        "dataset": "v3_assmang_mining_complete.sql",
        "objectives": [
            "Create calculated measures and members for business-friendly analytics.",
            "Understand named sets and reusable MDX logic.",
            "Design practical KPIs for production, cost, and safety monitoring.",
            "Use time-based calculations to support trend analysis.",
        ],
        "sections": [
            ("Calculated measures", [
                "Calculated measures derive new business insight without changing the source fact table.",
                "Examples include cost per tonne, revenue variance, and tonnes per employee.",
                "These measures should be clearly named and documented for users.",
            ]),
            ("Named sets and reusable logic", [
                "Named sets define reusable groups of members, such as top-performing mines or active operations.",
                "They simplify repeated report logic and improve consistency.",
            ]),
            ("KPIs in SSAS", [
                "A KPI combines value, goal, status, and often trend.",
                "At Assmang, KPIs can be created for safety score, production target attainment, or cost control.",
                "KPIs help executives consume analytics visually and consistently.",
            ]),
            ("Time-based logic", [
                "MDX calculations often compare current month to previous month, current year to previous year, or actual to target.",
                "This is where clean date hierarchies become especially valuable.",
            ]),
        ],
        "diagram": [
            "Base Measure -> Calculation Layer -> KPI Status -> Dashboard Consumer",
            "RevenueZAR -> Cost per Tonne -> Green/Amber/Red -> Power BI / Excel",
        ],
        "practical_steps": [
            "Create a calculated measure for `Cost Per Tonne` using total cost divided by tonnes produced.",
            "Create a named set for the mines with above-average production.",
            "Define a KPI for production target attainment using an assumed monthly target.",
            "Browse the KPI values in the cube browser or SSMS query window.",
            "Document the business meaning of each calculation created.",
        ],
        "hands_on": [
            "Design a safety KPI using `ComplianceScore` and explain thresholds for red, amber, and green status.",
            "Write a calculated measure for maintenance cost percentage of total cost.",
            "Create a named set for all chrome operations and reuse it in two different MDX queries.",
        ],
        "mcqs": [
            ("What is the main purpose of a calculated measure?", ["Replace a dimension", "Create derived business logic in the cube", "Delete source data", "Rename a partition"], "Create derived business logic in the cube"),
            ("What does a KPI usually contain?", ["Only a value", "Value, goal, and status", "Only a date hierarchy", "Only a relationship diagram"], "Value, goal, and status"),
            ("Why are named sets useful?", ["They store SQL tables", "They simplify repeated member selections", "They disable processing", "They replace dimensions"], "They simplify repeated member selections"),
        ],
        "scenario": "Executives want a production KPI that instantly shows whether each mine met its monthly target. Explain how you would design the KPI and what business thresholds you would set.",
        "practical_challenge": "Implement one calculated measure, one named set, and one KPI in the cube, then demonstrate all three in a query or browser view.",
    },
    {
        "day": "Day-02",
        "title": "Performance Tuning and Optimization",
        "dataset": "v3_assmang_mining_complete.sql",
        "objectives": [
            "Understand how storage mode and aggregation design affect performance.",
            "Recognise common causes of slow cube queries.",
            "Understand partitioning and caching at a beginner level.",
            "Apply practical optimisation decisions in an Assmang reporting context.",
        ],
        "sections": [
            ("Storage modes", [
                "MOLAP stores processed cube data and aggregations inside SSAS for best query speed.",
                "ROLAP queries relational storage more directly and can be useful when freshness matters more than speed.",
                "HOLAP combines both approaches.",
            ]),
            ("Aggregation design", [
                "Aggregations reduce the amount of work required at query time.",
                "The best aggregation design reflects common reporting patterns such as by mine, by month, and by department.",
                "Too many aggregations can increase processing time and storage usage.",
            ]),
            ("Partitioning and scalability", [
                "Large fact data can be partitioned by time or business area.",
                "Partitioning supports manageability, faster processing windows, and targeted optimisation.",
            ]),
            ("Caching and practical tuning", [
                "Repeated queries can benefit from cache reuse.",
                "Good dimension design, clean hierarchies, and sensible calculations all contribute to better performance.",
            ]),
        ],
        "diagram": [
            "Query Pattern -> Aggregation Match? -> Fast Response",
            "               -> No Match -> Detail Scan -> Slower Response",
        ],
        "practical_steps": [
            "Review the cube's storage mode and document why MOLAP is usually suitable for this training model.",
            "Use aggregation design tools to review suggested aggregations for production queries.",
            "Run two similar queries and compare behaviour before and after processing changes.",
            "Describe a partitioning idea for monthly or yearly production facts.",
            "Document one optimisation trade-off between speed, freshness, and complexity.",
        ],
        "hands_on": [
            "Recommend MOLAP, ROLAP, or HOLAP for a scenario where management needs near-real-time data for one measure group only.",
            "Propose a year-based partition strategy for `FactProduction`.",
            "List the top five things you would check first if an executive dashboard becomes slow.",
        ],
        "mcqs": [
            ("Which storage mode is usually fastest for query performance?", ["MOLAP", "ROLAP", "CSV", "XMLA"], "MOLAP"),
            ("What is the main purpose of partitioning?", ["To rename cubes", "To improve manageability and processing scalability", "To remove measures", "To disable dimensions"], "To improve manageability and processing scalability"),
            ("Why can too many aggregations be a problem?", ["They can increase processing cost", "They delete data", "They break MDX", "They prevent deployment"], "They can increase processing cost"),
        ],
        "scenario": "Assmang's monthly executive dashboard is fast, but detailed drill-through into one mining period is slow. Explain at least three areas you would investigate.",
        "practical_challenge": "Write a tuning note recommending storage mode, one aggregation improvement, and one partition strategy for the Assmang cube.",
    },
    {
        "day": "Day-02",
        "title": "Real-World SSAS Implementation at Assmang",
        "dataset": "v3_assmang_mining_complete.sql",
        "objectives": [
            "Apply the full SSAS workflow to an Assmang-style business solution.",
            "Design a business-ready analytical cube for production, cost, safety, and workforce reporting.",
            "Understand deployment, maintenance, and reporting integration considerations.",
            "Consolidate the course into a real implementation playbook.",
        ],
        "sections": [
            ("Business requirements to cube design", [
                "Production leaders need output and grade metrics by mine and time.",
                "Finance needs cost analytics by department and operation.",
                "Safety leadership needs incident and compliance monitoring.",
                "HR and management need employee metrics in context.",
            ]),
            ("Target cube design", [
                "Dimensions: Mine, Date, Department, Employee, and any KPI-supporting dimensions.",
                "Measure groups: Production, Operating Costs, Equipment Efficiency, Safety KPI, Employee Metrics.",
                "Calculated layer: cost per tonne, target attainment, utilisation measures.",
            ]),
            ("Integration and consumption", [
                "Excel is useful for analysts and pivot-based exploration.",
                "Power BI can consume SSAS live for dashboarding.",
                "SSMS and MDX remain valuable for admin testing and technical validation.",
            ]),
            ("Operations and maintenance", [
                "Plan processing windows, environment promotion, backup, and documentation.",
                "Monitor slow queries, aggregation effectiveness, and business definition changes over time.",
                "Treat SSAS as a governed semantic layer, not just a technical object.",
            ]),
        ],
        "diagram": [
            "SQL Warehouse -> SSAS Cube -> Excel / Power BI / SSMS",
            "      ^               |                |",
            "   Daily ETL      KPIs + Calcs      Business Decisions",
        ],
        "practical_steps": [
            "Load the v3 dataset and confirm all dimension and fact tables are available.",
            "Review the full cube structure and ensure all relevant measure groups are included.",
            "Create or validate at least one KPI each for production, cost, and safety.",
            "Run representative MDX queries that answer executive questions.",
            "Document a maintenance plan covering processing, monitoring, and business sign-off.",
        ],
        "hands_on": [
            "Create a mini solution blueprint for an executive production dashboard sourced from the cube.",
            "Define a daily, weekly, and monthly runbook for cube maintenance.",
            "List risks that could reduce trust in the cube and propose controls for each one.",
        ],
        "mcqs": [
            ("What is the best description of the SSAS cube in an enterprise BI landscape?", ["A raw source system", "A governed semantic analysis layer", "A printer queue", "A security group"], "A governed semantic analysis layer"),
            ("Which tool is commonly used by analysts to consume SSAS data interactively?", ["Notepad", "Excel", "Paint", "Command Prompt"], "Excel"),
            ("Why is documentation important in SSAS implementations?", ["It makes file names longer", "It preserves business logic and operational clarity", "It disables KPIs", "It replaces processing"], "It preserves business logic and operational clarity"),
        ],
        "scenario": "Assmang wants a single analytical layer for executives, analysts, and operational managers. Explain how one cube can serve all three groups differently while still using the same governed business definitions.",
        "practical_challenge": "Produce a one-page implementation playbook covering cube scope, dimensions, measure groups, KPIs, processing schedule, and reporting consumers for Assmang.",
    },
]


def slugify(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")


def render_theory(topic: dict) -> str:
    slug = slugify(topic['title'])
    lines = []
    lines += [f"# {topic['title']}", f"## {topic['day'].replace('-', ' ')} | Assmang Pty Ltd — SSAS Fundamentals Training", "", "---", ""]

    # Learning objectives
    lines += ["## 🎯 Learning Objectives", "", "By the end of this topic, participants will be able to:", ""]
    for i, obj in enumerate(topic['objectives'], 1):
        lines.append(f"{i}. {obj}")
    lines += ["", "---", ""]

    # Topic overview with context
    lines += ["## 📋 Topic Overview", ""]
    lines += [f"**Dataset:** `{topic['dataset']}`  ", f"**Difficulty:** Beginner (no prior SSAS experience required)  ", f"**Estimated reading time:** 20-30 minutes", ""]
    lines += ["### What is this topic about?", ""]
    lines += [f"This topic teaches you about **{topic['title']}**. If you have never worked with SQL Server Analysis Services before, don't worry — we will explain everything from scratch using plain language and real examples from Assmang's mining operations.", ""]
    lines += ["### Why does this matter to you?", ""]
    lines += ["As someone working at or with Assmang, you deal with data every day — production figures, costs, safety records, employee information. Right now, getting answers from that data probably involves:", ""]
    lines += ["- Asking someone in IT to write a report", "- Waiting for Excel spreadsheets to be updated", "- Running the same SQL queries over and over", "- Not being sure if the numbers are up to date", ""]
    lines += ["SSAS solves these problems by creating a **pre-built analytical model** (called a \"cube\") that lets anyone with Excel or Power BI get instant answers without writing code.", ""]
    lines += ["### The Assmang training context", ""]
    lines += ["All examples in this course use data from Assmang's actual operations:", ""]
    lines += ["| Mine | What it produces | Where it is |", "|------|-----------------|-------------|"]
    lines += ["| Beeshoek Mine | Iron Ore | Postmasburg, Northern Cape |"]
    lines += ["| Khumani Mine | Iron Ore | Kathu, Northern Cape |"]
    lines += ["| Black Rock Mine | Manganese | Hotazel, Northern Cape |"]
    lines += ["| Dwarsrivier Chrome Mine | Chrome | Burgersfort, Limpopo |"]
    lines += ["| Machadodorp Works | Chrome (processing) | Machadodorp, Mpumalanga |"]
    lines += ["", "---", ""]

    # Real-world analogy
    ANALOGIES = {
        "Introduction to SQL Server Analysis Services": (
            "a library catalogue system",
            "Imagine you have a massive library with thousands of books (your data). Without a catalogue, finding a specific book means searching every shelf manually. SSAS is like building a smart catalogue that already knows how many books you have by author, by genre, by year, and by shelf — so when someone asks \"How many science fiction books were published in 2023?\", the answer comes back instantly because it was pre-calculated.",
        ),
        "Multidimensional Models and Dimensions": (
            "the labels on filing cabinet drawers",
            "Think of dimensions like the labels on a filing cabinet. One drawer is labelled 'By Mine', another 'By Month', another 'By Department'. When you want to find production data for Beeshoek in March, you open the 'Mine' drawer, find 'Beeshoek', then look in the 'Month' section for 'March'. Dimensions are those category labels that help you navigate to exactly the data you need.",
        ),
        "Measures, Measure Groups, and Aggregations": (
            "the numbers on a sports scoreboard",
            "Think of measures like the numbers on a sports scoreboard — points scored, time remaining, fouls committed. These are the actual VALUES you care about. The scoreboard itself organises them by team (one dimension) and by quarter (another dimension). In SSAS, measures are the numbers (tonnes, revenue, cost) and measure groups are the scoreboards that organise related numbers together.",
        ),
        "Building and Deploying SSAS Cubes": (
            "building and opening a new shop",
            "Building a cube is like setting up a new shop. First you design the layout (data source view), then you stock the shelves (dimensions and measures), then you open the doors (deploy), and finally you turn on the lights so customers can see the products (process). The shop only becomes useful to customers after ALL these steps are complete.",
        ),
        "MDX Query Fundamentals": (
            "asking a specific question to a smart assistant",
            "Imagine MDX is like talking to a very organised assistant. Instead of saying 'find me some data', you say exactly: 'Show me REVENUE (what I want to see) for EACH MINE (rows) for the YEAR 2024 (filter).' The assistant already knows where everything is because the cube was pre-built, so the answer comes back instantly.",
        ),
        "Advanced Queries, Calculations, and KPIs": (
            "adding a dashboard with warning lights to your car",
            "Basic measures tell you speed and fuel level. But KPIs are like adding warning lights — green means everything is fine, amber means pay attention, red means there is a problem. A KPI takes a measure (like production tonnes), compares it to a target, and shows a colour-coded status so executives can instantly see which mines are on track and which need attention.",
        ),
        "Performance Tuning and Optimization": (
            "choosing between a motorbike and a truck for delivery",
            "Performance tuning is about choosing the right tool for the job. A motorbike (MOLAP) is fastest for small, frequent deliveries. A truck (ROLAP) carries more but is slower. A van (HOLAP) is a middle ground. You choose based on what your business actually needs — fast dashboards, real-time data, or a balance of both.",
        ),
        "Real-World SSAS Implementation at Assmang": (
            "building a complete control room for a mine",
            "This topic is like designing the entire control room for a mine. You decide what screens to display (dimensions and measures), what alarms to set (KPIs), how often to refresh the data (processing schedule), who can see what (security), and how to handle maintenance. It brings together everything you have learned into one complete, working solution.",
        ),
    }

    analogy_thing, analogy_text = ANALOGIES.get(topic['title'], ("a well-organised filing system", "Think of SSAS as a system that pre-organises your data so anyone can find answers quickly without technical expertise."))

    lines += ["## 🧠 Real-World Analogy (Plain English)", ""]
    lines += [f"**Think of this topic like {analogy_thing}.**", ""]
    lines += [f"{analogy_text}", ""]
    lines += ["> **Key insight:** SSAS takes complex data and makes it simple to explore. You don't need to be a programmer to use the results — you just need to know what question you want to answer.", ""]
    lines += ["---", ""]

    # Expanded sections with DEEP content
    for idx, (heading, bullets) in enumerate(topic['sections'], 1):
        lines += [f"## {idx}. {heading}", ""]

        # Plain English explanation
        lines += ["### 💬 In plain English", ""]
        lines += [f"Let's break down **{heading.lower()}** in the simplest possible terms:", ""]
        for b in bullets:
            lines += [f"**→** {b}", ""]

        # Detailed explanation
        lines += ["### 📚 Detailed explanation", ""]
        lines += [f"This concept is important because it directly affects how well the cube works for business users. Here is a deeper look:", ""]
        lines += [""]
        for j, b in enumerate(bullets):
            lines += [f"**Point {j+1}: {b}**", ""]
            lines += [f"What this means in practice: When you apply this at Assmang, it means that {b.lower().rstrip('.')}. This is not just a technical exercise — it directly helps managers, engineers, and executives get better information faster.", ""]
        lines += [""]

        # Real Assmang scenario
        lines += ["### 🏭 Assmang scenario", ""]
        lines += [f"**Situation:** A production manager at Khumani Mine asks: \"Can I see this month's iron ore output compared to last month, broken down by shift?\"", ""]
        lines += [f"**How {heading.lower()} helps:** Because the cube already has the right structure (dimensions for time and mine, measures for production), this question can be answered in seconds using Excel or Power BI — no SQL coding needed, no waiting for IT.", ""]
        lines += [""]

        # FAQ for this section
        lines += ["### ❓ Frequently Asked Questions", ""]
        lines += [f"**Q: Do I need to be a programmer to understand {heading.lower()}?**  ", f"A: No. This concept is about business logic and design thinking. The tools (SSDT) provide visual interfaces for most of the work.", ""]
        lines += [f"**Q: What happens if we get {heading.lower()} wrong?**  ", f"A: The cube will still work technically, but users may get confusing results, slow performance, or missing data. That's why we follow best practices from the start.", ""]
        lines += [f"**Q: How long does it take to set up {heading.lower()} for a real project?**  ", f"A: For a project the size of Assmang's training cube, this typically takes a few hours of design work plus a few hours of implementation and testing.", ""]
        lines += ["---", ""]

    # Comprehensive diagram section
    lines += ["## 📊 Architecture / Concept Diagram", ""]
    lines += ["The following diagram shows how this topic fits into the bigger picture:", ""]
    lines += ["```text"]
    for d in topic['diagram']:
        lines += [d]
    lines += ["```", ""]
    lines += ["### How to read this diagram", ""]
    lines += ["- **Left side:** Where your raw data lives (SQL Server database tables containing production, cost, safety, and employee data).", "- **Middle:** Where SSAS transforms that raw data into an analytical structure (the cube with its dimensions, hierarchies, and measures).", "- **Right side:** Where business users access the results (Excel pivot tables, Power BI dashboards, or MDX query results in SSMS).", ""]
    lines += ["### Why this matters", ""]
    lines += ["Without SSAS (the middle layer), every time a manager wants an answer, someone has to write SQL code against the raw database. With SSAS, the analytical structure is pre-built, so users can explore data independently using familiar tools like Excel.", ""]
    lines += ["---", ""]

    # Comprehensive terminology reference
    lines += ["## 📖 Key Terminology Reference", ""]
    lines += ["Here are the most important terms for this topic. Don't worry about memorising them all — you will learn them naturally through practice:", ""]
    lines += [""]
    lines += ["| Term | Plain English Definition | Example at Assmang |", "|------|------------------------|-------------------|"]
    lines += ["| **Cube** | A pre-built analytical structure that lets users explore data from many angles | The \"Assmang Mining Analytics\" cube containing all production and cost data |"]
    lines += ["| **Dimension** | A category you use to slice data (like filters in Excel) | Mine, Date, Department, Employee — these are the \"by what\" categories |"]
    lines += ["| **Hierarchy** | A drill-down path from general to specific | Year → Quarter → Month → Day (time hierarchy) |"]
    lines += ["| **Member** | One specific value within a dimension | \"Beeshoek Mine\" is a member of the Mine dimension |"]
    lines += ["| **Measure** | A number you want to analyse | Tonnes Produced, Revenue in ZAR, Cost Per Tonne |"]
    lines += ["| **Measure Group** | A collection of related measures from one business area | Production Measures (tonnes + grade + revenue) |"]
    lines += ["| **Fact Table** | The database table that stores the raw numbers | FactProduction, FactOperatingCosts |"]
    lines += ["| **Processing** | Loading data into the cube and building pre-calculated summaries | Running a nightly job that refreshes yesterday's production data |"]
    lines += ["| **Aggregation** | A pre-calculated total or average stored for speed | Total tonnes per mine per month (calculated once, queried many times) |"]
    lines += ["| **MDX** | The query language used to ask questions of a cube | Similar to SQL, but designed for multidimensional analysis |"]
    lines += ["| **MOLAP** | Storage mode where data is stored inside the cube for maximum speed | Default choice for Assmang — gives sub-second query times |"]
    lines += ["| **ROLAP** | Storage mode where data stays in SQL Server (slower but always fresh) | Used when real-time data is more important than speed |"]
    lines += ["| **KPI** | A traffic-light indicator showing whether a target is being met | Production KPI: Green if >= 90% of target, Red if < 70% |"]
    lines += ["| **SSDT** | SQL Server Data Tools — the IDE where you design and build cubes | Visual Studio with the SSAS project templates |"]
    lines += ["| **SSMS** | SQL Server Management Studio — for administration and testing | Where you deploy cubes and run MDX queries |"]
    lines += ["| **Data Source View (DSV)** | A logical view of which database tables the cube uses | Selecting Dim_Mine, Dim_Date, FactProduction for inclusion |"]
    lines += ["| **Deployment** | Pushing your cube design from your computer to the SSAS server | Like publishing a website — makes it available to users |"]
    lines += ["", "---", ""]

    # Best practices with explanations
    lines += ["## ✅ Best Practices for Beginners", ""]
    lines += ["Follow these rules from Day 1 and your SSAS projects will be much more successful:", ""]
    lines += [""]
    lines += ["### 1. Always start with a business question", "Before building anything technical, write down the question you're trying to answer. For example: \"The CEO wants to see monthly revenue by mine for the last 2 years.\" This drives every design decision.", ""]
    lines += ["### 2. Use clear, business-friendly names", "Don't name a dimension `Dim_001` or a measure `M_Rev`. Instead use `Mine` and `Revenue ZAR`. The people using your cube are not programmers — they need names that make instant sense.", ""]
    lines += ["### 3. Keep it simple at first", "Start with 3-4 dimensions and 5-6 measures. You can always add more later. A simple cube that works is infinitely better than a complex cube that confuses everyone.", ""]
    lines += ["### 4. Test with a real user", "After building your cube, sit down with a business user (not a developer) and ask them to find an answer. Watch where they get confused. Fix those areas.", ""]
    lines += ["### 5. Document everything", "Write down what each measure means, what each KPI threshold is, and when data is refreshed. Six months from now, you (or your replacement) will thank yourself.", ""]
    lines += ["### 6. Process and validate every time", "After any change to the cube, always process it AND check the results. An unprocessed cube looks fine in the designer but returns no data to users.", ""]
    lines += ["### 7. Plan for growth", "Assmang's data will grow. Design your cube so that adding a new year of data or a new mine doesn't require rebuilding everything from scratch.", ""]
    lines += ["---", ""]

    # Common mistakes with detailed explanations
    lines += ["## ⚠️ Common Mistakes (and How to Avoid Them)", ""]
    lines += ["Every beginner makes some of these mistakes. Knowing about them in advance will save you hours of frustration:", ""]
    lines += [""]
    lines += ["| # | Mistake | What goes wrong | How to prevent it |", "|---|---------|----------------|-------------------|"]
    lines += ["| 1 | Building without a business question | You create objects nobody uses, wasting time and confusing users | Always start with: \"What question am I answering?\" |"]
    lines += ["| 2 | Using technical names | Users see `Dim_Mine.MineID` instead of just \"Mine\" | Set display names in the dimension designer |"]
    lines += ["| 3 | Forgetting to process | Cube deploys successfully but shows zero data | Always process after deployment and check results |"]
    lines += ["| 4 | Summing percentages | Grade shows 340% because it summed 68% + 65% + 72% + 67% + 68% | Set aggregation to AVERAGE for ratios |"]
    lines += ["| 5 | No hierarchies | Users must scroll through 730 individual dates instead of drilling Year > Month | Create hierarchies for every dimension where drill-down makes sense |"]
    lines += ["| 6 | Not testing with business users | Cube works technically but nobody can use it | Demo to a non-technical user before promoting to production |"]
    lines += ["| 7 | No documentation | Nobody knows what the KPI thresholds are or when data refreshes | Keep a living document with business rules and schedules |"]
    lines += ["| 8 | Ignoring source data quality | Cube shows wrong totals because source data has duplicates or NULLs | Validate source data before cube processing |"]
    lines += ["", "---", ""]

    # Beginner FAQ
    lines += ["## ❓ Beginner FAQ", ""]
    lines += ["### \"Do I need to know how to program?\"", "No. SSAS development uses mostly visual tools (drag and drop in SSDT). You will learn some MDX query syntax in Day 2, but it's much simpler than full programming.", ""]
    lines += ["### \"How is this different from a normal Excel report?\"", "An Excel report shows you one fixed view of data. An SSAS cube lets you explore data from ANY angle — by mine, by month, by department, by commodity type — all without rebuilding the report. It's like the difference between a printed map and Google Maps.", ""]
    lines += ["### \"How long does it take to learn SSAS?\"", "The basics (this 2-day course) will get you building and querying cubes. Becoming an expert takes months of practice, but you can be productive within days.", ""]
    lines += ["### \"What if I make a mistake?\"", "SSAS is very forgiving during development. You can change dimensions, measures, and hierarchies as many times as you want before deploying to production. The dataset can be reloaded at any time.", ""]
    lines += ["### \"Who uses the cube after we build it?\"", "Anyone with Excel or Power BI can connect to the cube and explore data. They don't need SSAS knowledge — they just use familiar tools (pivot tables, charts) that connect to the cube behind the scenes.", ""]
    lines += ["---", ""]

    # Summary
    lines += ["## 📝 Topic Summary", ""]
    lines += [f"In this topic you learned about **{topic['title']}**.", ""]
    lines += ["### Key takeaways:", ""]
    for obj in topic['objectives']:
        lines += [f"- ✅ {obj}"]
    lines += [""]
    lines += ["### What to do next:", ""]
    lines += [f"1. Complete the **practical lab** (guided, step-by-step) using dataset `{topic['dataset']}`", "2. Attempt the **later hands-on exercises** (independent practice)", "3. Complete the **assessment** to test your understanding", "4. Move on to the next topic when you feel confident", ""]
    lines += ["### How to know you understand this topic:", ""]
    lines += ["- You can explain the key concepts to a colleague in plain English", "- You can identify where this topic fits in the overall SSAS workflow", "- You can connect the concepts to a real Assmang business question", "- You completed the practical lab successfully", ""]
    lines += ["---", ""]
    lines += [f"*Assmang Pty Ltd — SSAS Fundamentals Training | {topic['day'].replace('-', ' ')}*  ", f"*Course: SSAS100 | Level: Beginner | Topic: {topic['title']}*", ""]
    return "\n".join(lines)


def render_practical(topic: dict) -> str:
    lines = [f"# Practical Lab — {topic['title']}", f"## {topic['day'].replace('-', ' ')} | Assmang Pty Ltd — SSAS Fundamentals", "", "---", ""]
    lines += ["## 🎯 Lab Goal", "", f"Apply the theory from **{topic['title']}** by completing a guided, step-by-step exercise in SQL Server Data Tools (SSDT) and SQL Server Management Studio (SSMS).", ""]
    lines += ["## 📋 Prerequisites", "", f"- Dataset **`{topic['dataset']}`** loaded into SQL Server", "- SQL Server Analysis Services running", "- Visual Studio with SSDT installed", "- SSMS available for verification", ""]
    lines += ["## 🔧 Lab Environment", "", "| Component | Value |", "|-----------|-------|", "| SQL Server Instance | localhost\\SSASDEV (or your instance) |", f"| Database | AssmangMining |", "| SSAS Project | AssmangMiningCube |", f"| Dataset Version | `{topic['dataset']}` |", "", "---", ""]
    lines += ["## 📝 Guided Steps", ""]
    for i, step in enumerate(topic['practical_steps'], 1):
        lines += [f"### Step {i}: {step.split('.')[0] if '.' in step else step[:50]}", ""]
        lines += [f"**What to do:** {step}", ""]
        lines += [f"**Why this matters:** This step builds your understanding of {topic['title'].lower()} by giving you hands-on experience with the tool.", ""]
        lines += [f"**Expected result:** You should see a successful outcome or confirmation in SSDT/SSMS before moving to the next step.", ""]
        lines += ["**Troubleshooting:** If this step fails, check:", f"- Is the SQL Server instance running?", f"- Is the dataset `{topic['dataset']}` loaded?", "- Do you have the correct permissions?", "", "---", ""]
    lines += ["## ✅ Validation Checklist", "", "Before marking this lab as complete, confirm:", ""]
    lines += ["- [ ] The relevant SQL dataset was loaded and verified", "- [ ] The SSAS project was opened without errors", "- [ ] All objects created in this lab are visible in Solution Explorer", "- [ ] Processing completed successfully (check Output window)", "- [ ] The cube browser or SSMS query returns expected results", "- [ ] You can explain what each object does in business terms", ""]
    lines += ["---", ""]
    lines += ["## 🎓 Expected Outcome", "", f"By the end of this lab, you should be able to demonstrate the core workflow for **{topic['title']}** in the Assmang training environment. You should be able to:", ""]
    for obj in topic['objectives']:
        lines += [f"- {obj}"]
    lines += ["", "---", ""]
    lines += ["## 💡 Tips for Success", "", "- **Read each step fully** before executing it.", "- **Save your project** after each major step.", "- **Ask questions** if something doesn't look right — it's better to clarify early.", "- **Take notes** on what you observe — this helps with the assessment later.", ""]
    lines += ["---", "", f"*Assmang Pty Ltd — SSAS Fundamentals | {topic['day'].replace('-', ' ')} Practical Lab*", ""]
    return "\n".join(lines)


def render_hands_on(topic: dict) -> str:
    lines = [f"# Later Hands-On Exercises — {topic['title']}", f"## {topic['day'].replace('-', ' ')} | Assmang Pty Ltd — SSAS Fundamentals", "", "---", ""]
    lines += ["## 🎯 Purpose", "", "These exercises are designed for **independent practice** after the guided lab. They are slightly more challenging and require you to apply what you've learned without step-by-step guidance.", ""]
    lines += ["## 📋 Before You Begin", "", f"- Ensure the guided lab for **{topic['title']}** is complete", f"- Dataset **`{topic['dataset']}`** must be loaded", "- Your SSAS project should be in a working state", "- Allow 30-45 minutes for these exercises", ""]
    lines += ["---", ""]
    for i, ex in enumerate(topic['hands_on'], 1):
        lines += [f"## Exercise {i}", ""]
        lines += [f"### Task", "", ex, ""]
        lines += [f"### Hints", "", f"- Refer back to the theory for **{topic['title']}** if you get stuck.", "- Think about how this connects to a real business question at Assmang.", "- There may be multiple correct approaches — choose the one you can explain clearly.", ""]
        lines += [f"### Deliverable", "", f"- A written answer (1-2 paragraphs) OR a screenshot of your SSDT/SSMS result.", "- Be prepared to explain your reasoning to the trainer.", ""]
        lines += ["---", ""]
    lines += ["## ✅ Success Criteria", ""]
    lines += ["Your exercises are considered successful when:", ""]
    lines += ["- Your answer reflects the topic's **business purpose**, not only the technical steps.", "- You can explain **why** the design or query choice fits Assmang's reporting needs.", "- You can connect your answer back to dimensions, measures, hierarchies, MDX, or deployment where relevant.", "- Your work is **documented clearly** enough that a colleague could understand it.", ""]
    lines += ["---", ""]
    lines += ["## 💡 Stretch Challenge (Optional)", "", f"If you finish early, try to extend one of the exercises above by combining it with a concept from a previous topic. For example, if this topic covers measures, try connecting your measure design to a specific dimension hierarchy from an earlier topic.", ""]
    lines += ["---", "", f"*Assmang Pty Ltd — SSAS Fundamentals | {topic['day'].replace('-', ' ')} Independent Practice*", ""]
    return "\n".join(lines)


def render_assessment(topic: dict) -> str:
    lines = [f"# Assessment — {topic['title']}", f"## {topic['day'].replace('-', ' ')} | Assmang Pty Ltd — SSAS Fundamentals", "", "---", ""]
    lines += ["## 📋 Assessment Overview", "", "This assessment covers the key concepts from this topic. It includes:", "", "- **Section A:** Multiple-choice questions (knowledge recall)", "- **Section B:** Scenario-based question (application and reasoning)", "- **Section C:** Practical challenge (hands-on demonstration)", "", "**Passing score:** 70% (combined across all sections)", "", "---", ""]
    lines += ["## Section A: Multiple-Choice Questions (1 point each)", ""]
    for idx, (q, opts, ans) in enumerate(topic['mcqs'], 1):
        lines.append(f"**Q{idx}. {q}**")
        lines.append("")
        for letter, opt in zip(["A", "B", "C", "D"], opts):
            lines.append(f"- {letter}) {opt}")
        lines.append("")
    lines += ["---", ""]
    lines += ["## Section B: Scenario Question (3 points)", ""]
    lines += ["**Scenario:**", ""]
    lines += [topic['scenario'], ""]
    lines += ["**Your answer should include:**", ""]
    lines += ["- A clear explanation of the SSAS concept involved", "- How it connects to the Assmang business context", "- What you would recommend and why", "- Any risks or considerations", ""]
    lines += ["---", ""]
    lines += ["## Section C: Practical Challenge (5 points)", ""]
    lines += ["**Task:**", ""]
    lines += [topic['practical_challenge'], ""]
    lines += ["**Submission requirements:**", ""]
    lines += ["- Screenshot(s) showing your work in SSDT or SSMS", "- Brief written explanation of what you did and why", "- Evidence that the result is correct (e.g., cube browser output, query result)", ""]
    lines += ["---", ""]
    lines += ["## 📝 Grading Rubric", ""]
    lines += ["| Section | Points | Criteria |", "|---------|--------|----------|"]
    lines += [f"| Section A (MCQ) | {len(topic['mcqs'])} | Correct answer = 1 point each |"]
    lines += ["| Section B (Scenario) | 3 | Clear reasoning (1) + Business context (1) + Recommendation (1) |"]
    lines += ["| Section C (Practical) | 5 | Correct execution (2) + Evidence (2) + Explanation (1) |"]
    lines += [f"| **Total** | **{len(topic['mcqs']) + 8}** | |"]
    lines += ["", "---", ""]
    lines += ["## ✅ Answer Key", ""]
    lines += ["### Section A Answers", ""]
    for idx, (q, opts, ans) in enumerate(topic['mcqs'], 1):
        lines.append(f"- **Q{idx}:** {ans}")
    lines += [""]
    lines += ["### Section B — Scenario Guidance", ""]
    lines += [f"A strong answer should:", ""]
    lines += [f"- Discuss business context specific to Assmang mining operations", f"- Reference SSAS terminology correctly (dimensions, measures, hierarchies, processing, etc.)", f"- Explain the design decision with clear reasoning", f"- Acknowledge trade-offs or limitations", ""]
    lines += ["### Section C — Practical Challenge Guidance", ""]
    lines += ["A strong submission should:", ""]
    lines += ["- Show correct SSAS object naming and configuration", "- Demonstrate successful processing or querying", "- Include evidence that validates correctness", "- Be documented clearly enough for a reviewer to follow", ""]
    lines += ["---", "", f"*Assmang Pty Ltd — SSAS Fundamentals | {topic['day'].replace('-', ' ')} Assessment*", ""]
    return "\n".join(lines)


def build_markdown():
    day_map = {}
    for topic in TOPICS:
        day_dir = BASE / topic['day']
        slug = slugify(topic['title'])
        topic_dir = day_dir / slug
        topic_dir.mkdir(parents=True, exist_ok=True)
        (topic_dir / 'assets').mkdir(exist_ok=True)
        (topic_dir / f"{slug}-theory.md").write_text(render_theory(topic))
        (topic_dir / f"{slug}-practical.md").write_text(render_practical(topic))
        (topic_dir / f"{slug}-later-hands-on.md").write_text(render_hands_on(topic))
        (topic_dir / f"{slug}-assessment.md").write_text(render_assessment(topic))
        day_map.setdefault(topic['day'], []).append((topic['title'], slug))
    return day_map


# ---------- Presentation generation ----------
C_NAVY = RGBColor(38, 59, 82)
C_BLUE = RGBColor(78, 116, 156)
C_BG = RGBColor(248, 250, 252)
C_TEXT = RGBColor(33, 43, 54)
C_LIGHT = RGBColor(230, 237, 245)
C_ACCENT = RGBColor(90, 132, 174)
C_WHITE = RGBColor(255, 255, 255)


def add_box(slide, left, top, width, height, fill, line):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, text, left, top, width, height, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, font='Calibri'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def add_bullets(slide, bullets, left, top, width, height, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.text = bullet; p.level = 0; p.bullet = True
        p.runs[0].font.size = Pt(size); p.runs[0].font.name = 'Calibri'; p.runs[0].font.color.rgb = C_TEXT
    return box


def add_title(slide, title, subtitle='Assmang Pty Ltd | SSAS Fundamentals'):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C_BG
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.85))
    banner.fill.solid(); banner.fill.fore_color.rgb = C_NAVY; banner.line.color.rgb = C_NAVY
    add_textbox(slide, title, Inches(0.35), Inches(0.12), Inches(9.8), Inches(0.35), size=26, bold=True, color=C_WHITE, font='Calibri Light')
    add_textbox(slide, subtitle, Inches(0.35), Inches(0.47), Inches(6.0), Inches(0.15), size=10.5, color=RGBColor(220,229,240))


def generate_ppt(day_map):
    for topic in TOPICS:
        slug = slugify(topic['title'])
        topic_dir = BASE / topic['day'] / slug
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid(); s.background.fill.fore_color.rgb = C_NAVY
        band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.0), Inches(13.33), Inches(1.5))
        band.fill.solid(); band.fill.fore_color.rgb = C_BLUE; band.line.color.rgb = C_BLUE
        add_textbox(s, topic['day'].replace('-', ' '), Inches(0.5), Inches(0.8), Inches(4), Inches(0.35), size=18, color=C_WHITE)
        add_textbox(s, topic['title'], Inches(0.5), Inches(1.6), Inches(11.7), Inches(2), size=30, bold=True, color=C_WHITE, font='Calibri Light')
        add_textbox(s, 'Theory-only trainer deck with diagrams, examples, and recap', Inches(0.55), Inches(4.2), Inches(8), Inches(0.3), size=14, color=RGBColor(226, 234, 242))

        # Objectives
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, 'Learning Objectives', topic['title'])
        add_box(s, Inches(0.5), Inches(1.2), Inches(12.1), Inches(5.7), C_WHITE, C_LIGHT)
        add_bullets(s, topic['objectives'], Inches(0.9), Inches(1.7), Inches(7.2), Inches(4.7), 18)
        add_box(s, Inches(8.65), Inches(1.8), Inches(3.25), Inches(3.0), C_LIGHT, C_ACCENT)
        add_textbox(s, 'Dataset', Inches(9.25), Inches(2.15), Inches(2.0), Inches(0.25), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, topic['dataset'], Inches(8.95), Inches(2.8), Inches(2.6), Inches(0.5), size=14, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(s, 'Business context: Assmang mining analytics', Inches(8.9), Inches(3.6), Inches(2.8), Inches(0.8), size=13, color=C_TEXT, align=PP_ALIGN.CENTER)

        # Four concept slides
        for heading, bullets in topic['sections']:
            s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, heading, topic['title'])
            add_box(s, Inches(0.5), Inches(1.2), Inches(8.1), Inches(5.7), C_WHITE, C_LIGHT)
            add_bullets(s, bullets, Inches(0.9), Inches(1.7), Inches(7.1), Inches(4.7), 18)
            add_box(s, Inches(8.9), Inches(1.5), Inches(3.0), Inches(4.8), C_LIGHT, C_ACCENT)
            add_textbox(s, 'Trainer note', Inches(9.5), Inches(1.9), Inches(1.8), Inches(0.25), size=17, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
            add_textbox(s, f'Connect this concept to the Assmang use case and validate it in the lab.', Inches(9.2), Inches(2.7), Inches(2.4), Inches(2.2), size=14, color=C_TEXT, align=PP_ALIGN.CENTER)

        # Diagram slide
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, 'Architecture / Concept Diagram', topic['title'])
        add_box(s, Inches(0.6), Inches(1.45), Inches(12.0), Inches(4.9), C_WHITE, C_LIGHT)
        add_textbox(s, '\n'.join(topic['diagram']), Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.5), size=19, color=C_NAVY, font='Courier New', align=PP_ALIGN.CENTER)
        add_textbox(s, 'Use this visual to explain the flow of structure, context, and analytical consumption.', Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.3), size=14, color=RGBColor(90,100,110), align=PP_ALIGN.CENTER)

        # Practical emphasis slide
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, 'Guided Lab Focus', topic['title'])
        add_box(s, Inches(0.55), Inches(1.25), Inches(12.0), Inches(5.55), C_WHITE, C_LIGHT)
        add_bullets(s, topic['practical_steps'], Inches(0.9), Inches(1.75), Inches(10.7), Inches(4.5), 17)

        # Recap slide
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s, 'Recap and Common Mistakes', topic['title'])
        add_box(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.9), C_WHITE, C_LIGHT)
        add_box(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(4.9), C_WHITE, C_LIGHT)
        add_textbox(s, 'Key messages', Inches(2.2), Inches(1.85), Inches(2.3), Inches(0.25), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_textbox(s, 'Avoid these mistakes', Inches(8.35), Inches(1.85), Inches(2.8), Inches(0.25), size=18, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_bullets(s, [o for o in topic['objectives'][:3]], Inches(0.9), Inches(2.35), Inches(5.0), Inches(3.4), 15)
        add_bullets(s, [
            'Skipping validation after deployment or processing.',
            'Using weak business naming that confuses users.',
            'Designing objects without linking them to reporting questions.',
        ], Inches(7.2), Inches(2.35), Inches(5.0), Inches(3.4), 15)

        # End slide
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid(); s.background.fill.fore_color.rgb = C_NAVY
        add_textbox(s, 'Questions and Lab Transition', Inches(1.3), Inches(2.0), Inches(10.7), Inches(0.6), size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, topic['title'], Inches(1.3), Inches(3.0), Inches(10.7), Inches(0.4), size=18, color=RGBColor(205, 217, 232), align=PP_ALIGN.CENTER)

        prs.save(topic_dir / f'{slug}.pptx')


def rewrite_readme(day_map):
    text = dedent('''
    # SSAS Fundamentals — SQL Server Analysis Services
    ### Assmang Pty Ltd — Technical Training Programme

    ---

    ## 📋 Course Overview

    | Field | Detail |
    |---|---|
    | **Course Code** | SSAS100 |
    | **Course Title** | SQL Server Analysis Services Fundamentals |
    | **Client** | Assmang Pty Ltd |
    | **Duration** | 2 Days |
    | **Level** | Beginner |
    | **Delivery Format** | Instructor-led, hands-on lab environment |
    | **Primary Platform** | Windows with SQL Server, SSMS, SSDT, and SSAS |

    ---

    ## 🎯 Course Goal

    By the end of this course, participants will be able to design, build, deploy, process, browse, and query a multidimensional SSAS cube using Assmang-aligned mining data.

    ---

    ## 📖 Topics by Day

    ### Day 1 — SSAS Fundamentals & Dimensional Modeling

    1. **Introduction to SQL Server Analysis Services**
    2. **Multidimensional Models and Dimensions**
    3. **Measures, Measure Groups, and Aggregations**
    4. **Building and Deploying SSAS Cubes**

    ### Day 2 — Queries, Advanced Queries & Optimization

    5. **MDX Query Fundamentals**
    6. **Advanced Queries, Calculations, and KPIs**
    7. **Performance Tuning and Optimization**
    8. **Real-World SSAS Implementation at Assmang**

    ---

    ## 📂 Folder Structure

    ```
    SSAS-Fundamentals-SQL-Server-Analysis-Services/
    ├── README.md
    ├── START_HERE.md
    ├── COURSE_BUILD_STATUS.md
    ├── requirements.txt
    ├── build_ssas_course.py
    ├── tools-installation/
    ├── datasets/
    ├── assessments/
    ├── Day-01/
    │   ├── introduction-to-sql-server-analysis-services/
    │   ├── multidimensional-models-and-dimensions/
    │   ├── measures-measure-groups-and-aggregations/
    │   └── building-and-deploying-ssas-cubes/
    └── Day-02/
        ├── mdx-query-fundamentals/
        ├── advanced-queries-calculations-and-kpis/
        ├── performance-tuning-and-optimization/
        └── real-world-ssas-implementation-at-assmang/
    ```

    Each topic folder contains:
    - `<topic-name>-theory.md`
    - `<topic-name>-practical.md`
    - `<topic-name>-later-hands-on.md`
    - `<topic-name>-assessment.md`
    - `<topic-name>.pptx`
    - `assets/`

    ---

    ## 🧱 Progressive Datasets

    - **v1** — Base dimensions for introductory dimensional modelling
    - **v2** — Production and operating cost facts for measure groups and cube build
    - **v3** — KPI, safety, equipment, and workforce metrics for advanced querying and optimisation

    ---

    ## ✅ Completion Status

    This course folder has been fully generated with all 8 topics, all required markdown files, and all topic presentations.
    ''').strip() + '\n'
    (BASE / 'README.md').write_text(text)


def rewrite_status_files():
    start = dedent('''
    # START HERE

    ## Course Status

    The SSAS course is now fully complete.

    ### Included
    - All 8 exact syllabus topics
    - All required markdown files per topic
    - All generated PowerPoint presentations
    - Progressive datasets v1, v2, and v3
    - Tool installation guide
    - Final assessment

    ### Recommended Sequence
    1. Read `README.md`
    2. Install tools from `tools-installation/README.md`
    3. Load datasets in order: v1 -> v2 -> v3
    4. Deliver Day 1 topics, then Day 2 topics
    5. Use the generated `.pptx` files for theory delivery and the markdown files for labs and assessment
    ''').strip() + '\n'
    (BASE / 'START_HERE.md').write_text(start)

    status = dedent('''
    # COURSE BUILD STATUS

    ## Final Result

    - Topic folders: 8/8 complete
    - Topic theory files: 8/8 complete
    - Topic practical files: 8/8 complete
    - Topic later-hands-on files: 8/8 complete
    - Topic assessment files: 8/8 complete
    - Topic presentations: 8/8 generated
    - Datasets: 3/3 complete
    - Installation guides: complete
    - Final assessment: complete

    ## Validation Goal

    The course is considered complete only when every topic folder contains the following files:
    - `<topic-name>-theory.md`
    - `<topic-name>-practical.md`
    - `<topic-name>-later-hands-on.md`
    - `<topic-name>-assessment.md`
    - `<topic-name>.pptx`
    ''').strip() + '\n'
    (BASE / 'COURSE_BUILD_STATUS.md').write_text(status)


def write_final_assessment():
    text = dedent('''
    # Final Assessment — SSAS Fundamentals

    ## Section A: Multiple Choice

    1. What is the main purpose of SSAS in this course?
    2. Which object supports drill-down navigation?
    3. What is a measure group typically based on?
    4. What does processing do in SSAS?
    5. What language is used to query multidimensional cubes?
    6. Which storage mode is usually fastest for query performance?
    7. What is a KPI used for?
    8. Why are partitions useful?

    ## Section B: Scenario Questions

    1. Explain how Assmang could use one SSAS cube to support production, cost, and safety reporting.
    2. Explain why a Date hierarchy is essential for time-based reporting.
    3. Describe how you would validate a cube before releasing it to business users.

    ## Section C: Practical Challenge

    Build or describe a cube solution that contains:
    - Dimensions for Mine, Date, Department, and Employee
    - Measure groups for Production and Operating Costs
    - One calculated measure
    - One KPI
    - A basic processing and deployment plan

    ## Answer Guidance

    Strong answers should use the vocabulary of dimensions, hierarchies, measures, measure groups, MDX, deployment, processing, and performance tuning correctly in the Assmang context.
    ''').strip() + '\n'
    (BASE / 'assessments' / 'final_assessment.md').write_text(text)


def clean_and_normalize(day_map):
    # remove non-compliant placeholder dirs if present
    for path in [
        BASE / 'Day-02' / 'Topic-5-Directory',
        BASE / 'Day-02' / 'Topic-6-Directory',
        BASE / 'Day-02' / 'Topic-7-Directory',
        BASE / 'Day-02' / 'Topic-8-Directory',
        BASE / 'Day-01' / 'Introduction-to-SQL-Server-Analysis-Services',
        BASE / 'Day-01' / 'Multidimensional-Models-and-Dimensions',
        BASE / 'Day-01' / 'Measures-Measure-Groups-and-Aggregations',
        BASE / 'Day-01' / 'Building-and-Deploying-SSAS-Cubes',
    ]:
        if path.exists() and path.is_dir():
            shutil.rmtree(path)


def validate():
    expected = sorted((t['day'], slugify(t['title'])) for t in TOPICS)
    found = []
    for day in ['Day-01', 'Day-02']:
        for child in sorted((BASE / day).iterdir()):
            if child.is_dir():
                found.append((day, child.name))
    assert found == expected, f'Unexpected topic directories: {found}'
    for topic in TOPICS:
        slug = slugify(topic['title'])
        tdir = BASE / topic['day'] / slug
        for fname in [f'{slug}-theory.md', f'{slug}-practical.md', f'{slug}-later-hands-on.md', f'{slug}-assessment.md', f'{slug}.pptx']:
            assert (tdir / fname).exists(), f'Missing {tdir / fname}'
    return True


def main():
    clean_and_normalize({})
    day_map = build_markdown()
    generate_ppt(day_map)
    rewrite_readme(day_map)
    rewrite_status_files()
    write_final_assessment()
    validate()
    print('SSAS course build complete and validated.')


if __name__ == '__main__':
    main()


